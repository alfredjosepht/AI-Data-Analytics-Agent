from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
import os
from datetime import datetime

from backend.database.sqlite_manager import sqlite_manager
from backend.database.duckdb_manager import duckdb_manager

# PDF and PPTX packages import inside endpoint or globally
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception:
    letter = None
    SimpleDocTemplate = None
    Presentation = None

router = APIRouter(
    tags=["Export"]
)


@router.get("/workspaces/{workspace_id}/export/pdf")
def export_pdf(workspace_id: int):

    ws = sqlite_manager.get_workspace(workspace_id)
    if not ws:
        raise HTTPException(status_code=404, detail="Workspace not found")

    if SimpleDocTemplate is None:
        raise HTTPException(status_code=500, detail="reportlab not installed. Please wait for dependency installation to complete.")

    out_dir = "exports"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"workspace_{workspace_id}.pdf")

    doc = SimpleDocTemplate(out_path, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'TitleStyle',
        parent=styles['Heading1'],
        fontSize=24,
        leading=28,
        textColor=colors.HexColor('#1e293b'),
        spaceAfter=15
    )
    section_style = ParagraphStyle(
        'SectionStyle',
        parent=styles['Heading2'],
        fontSize=14,
        leading=18,
        textColor=colors.HexColor('#3b82f6'),
        spaceBefore=12,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'BodyStyle',
        parent=styles['Normal'],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#334155'),
        spaceAfter=6
    )

    elements = []
    
    elements.append(Paragraph(f"Enterprise Data Analytics Report", title_style))
    elements.append(Paragraph(f"<b>Workspace Name:</b> {ws.get('name').upper()}", body_style))
    elements.append(Paragraph(f"<b>Source File:</b> {ws.get('file_name')}", body_style))
    elements.append(Paragraph(f"<b>Generated At:</b> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", body_style))
    elements.append(Spacer(1, 15))

    # Dataset Summary Section
    elements.append(Paragraph("Dataset Overview & Metadata", section_style))
    meta = ws.get('metadata') or {}
    row_count = meta.get('row_count', 'N/A')
    col_count = meta.get('column_count', 'N/A')
    cols = meta.get('columns', [])
    elements.append(Paragraph(f"The dataset contains <b>{row_count}</b> rows and <b>{col_count}</b> columns.", body_style))
    elements.append(Paragraph(f"<b>Column Schema:</b> {', '.join(cols)}", body_style))
    elements.append(Spacer(1, 10))

    # Data Quality Summary
    elements.append(Paragraph("Data Quality Profile", section_style))
    cleaning_report = ws.get('cleaning_report') or {}
    report_desc = cleaning_report.get('report', 'No cleaning needed or applied.')
    approved = "Yes" if ws.get('cleaning_approved') else "No (Draft/Pending Confirm)"
    elements.append(Paragraph(f"<b>Cleaning Approved:</b> {approved}", body_style))
    elements.append(Paragraph(f"<b>Cleaning Steps Log:</b> {report_desc}", body_style))
    elements.append(Spacer(1, 10))

    # Queries & Insights
    chats = sqlite_manager.get_chat_history(workspace_id)
    if chats:
        elements.append(Paragraph("AI Analytical Insights & Recommendations", section_style))
        for chat in chats[-3:]:
            q_time = chat.get('created_at', '')
            elements.append(Paragraph(f"<b>Query (at {q_time}):</b> {chat.get('question')}", body_style))
            elements.append(Paragraph(f"<b>AI Executive Insight:</b> {chat.get('answer')}", body_style))
            elements.append(Spacer(1, 8))
    else:
        table_name = ws.get('table_name')
        if table_name:
            elements.append(Paragraph("Sample Dataset Rows (Top 10)", section_style))
            try:
                df = duckdb_manager.query(f"SELECT * FROM {table_name} LIMIT 10")
                data_list = [list(df.columns)]
                for idx, row in df.iterrows():
                    data_list.append([str(val)[:20] for val in row])
                
                col_widths = [letter[0] / (len(df.columns) + 1)] * len(df.columns)
                t = Table(data_list, colWidths=col_widths)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1e293b')),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                    ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                    ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0,0), (-1,-1), 8),
                    ('BOTTOMPADDING', (0,0), (-1,0), 4),
                    ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#f8fafc')),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
                    ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor('#f8fafc'), colors.HexColor('#ffffff')])
                ]))
                elements.append(t)
            except Exception as e:
                elements.append(Paragraph(f"Failed to query preview: {e}", body_style))

    doc.build(elements)
    return FileResponse(out_path, media_type="application/pdf", filename=os.path.basename(out_path))


@router.get("/workspaces/{workspace_id}/export/pptx")
def export_pptx(workspace_id: int):

    ws = sqlite_manager.get_workspace(workspace_id)
    if not ws:
        raise HTTPException(status_code=404, detail="Workspace not found")

    if Presentation is None:
        raise HTTPException(status_code=500, detail="python-pptx not installed. Please wait for dependency installation to complete.")

    out_dir = "exports"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"workspace_{workspace_id}.pptx")

    prs = Presentation()
    
    def apply_title(slide, text):
        title = slide.shapes.title
        title.text = text
        for p in title.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 41, 59)

    # 1. Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = ws.get('name').upper()
    slide1.placeholders[1].text = f"Automated AI Analytics Dashboard Report\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # 2. Slide 2: Dataset Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title(slide2, "Dataset Overview")
    body2 = slide2.shapes.placeholders[1].text_frame
    body2.text = f"Workspace Name: {ws.get('name')}"
    body2.add_paragraph().text = f"Source Dataset Filename: {ws.get('file_name')}"
    
    meta = ws.get('metadata') or {}
    body2.add_paragraph().text = f"Total Rows: {meta.get('row_count', 'N/A')}"
    body2.add_paragraph().text = f"Total Columns: {meta.get('column_count', 'N/A')}"
    cols = meta.get('columns', [])[:8]
    cols_str = ", ".join(cols) + ("..." if len(meta.get('columns', [])) > 8 else "")
    body2.add_paragraph().text = f"Columns parsed: {cols_str}"

    # 3. Slide 3: Data Quality & Cleaning
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title(slide3, "Data Quality & Operations Log")
    body3 = slide3.shapes.placeholders[1].text_frame
    cleaning_report = ws.get('cleaning_report') or {}
    report_desc = cleaning_report.get('report', 'No cleaning needed.')
    body3.text = "Operations executed automatically:"
    body3.add_paragraph().text = "• Whitespace trimmed and cases standardized."
    body3.add_paragraph().text = "• Empty rows and columns dropped."
    body3.add_paragraph().text = f"• Operations Log: {report_desc}"

    # 4. Slide 4: AI Insights
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title(slide4, "AI Executive Insights")
    body4 = slide4.shapes.placeholders[1].text_frame
    
    chats = sqlite_manager.get_chat_history(workspace_id)
    if chats:
        body4.text = "Key findings from user questions:"
        for i, chat in enumerate(chats[-2:]):
            body4.add_paragraph().text = f"Query {i+1}: '{chat.get('question')}'"
            p = body4.add_paragraph()
            p.text = f"Insight: {chat.get('answer')}"
            p.level = 1
    else:
        body4.text = "No questions run on this dataset yet."
        body4.add_paragraph().text = "Run natural language query questions on the chat panel to generate automated business insights."

    # 5. Slide 5: Recommendations & Next Steps
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title(slide5, "Recommendations & Next Steps")
    body5 = slide5.shapes.placeholders[1].text_frame
    body5.text = "Next Steps in Workspace Dashboard:"
    body5.add_paragraph().text = "1. Customize the Plotly interactive visualizations in the Center panel."
    body5.add_paragraph().text = "2. Filter, search, and sort data in the tabular view."
    body5.add_paragraph().text = "3. Create a Scheduled Report from the history page to run queries on cron."
    body5.add_paragraph().text = "4. Export refined PDFs or CSV sheets with your latest operations."

    prs.save(out_path)
    return FileResponse(out_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=os.path.basename(out_path))


@router.get("/workspaces/{workspace_id}/export/csv")
def export_csv(workspace_id: int):
    import pandas as pd
    ws = sqlite_manager.get_workspace(workspace_id)
    if not ws:
        raise HTTPException(status_code=404, detail="Workspace not found")
    table_name = ws.get("table_name")
    if not table_name:
        raise HTTPException(status_code=400, detail="No active table found for this workspace")
    try:
        df = duckdb_manager.query(f"SELECT * FROM {table_name}")
        out_dir = "exports"
        os.makedirs(out_dir, exist_ok=True)
        out_path = os.path.join(out_dir, f"workspace_{workspace_id}_cleaned.csv")
        df.to_csv(out_path, index=False)
        return FileResponse(out_path, media_type="text/csv", filename=f"{ws.get('name')}_export.csv")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/workspaces/{workspace_id}/export/excel")
def export_excel(workspace_id: int):
    import pandas as pd
    ws = sqlite_manager.get_workspace(workspace_id)
    if not ws:
        raise HTTPException(status_code=404, detail="Workspace not found")
    table_name = ws.get("table_name")
    if not table_name:
        raise HTTPException(status_code=400, detail="No active table found for this workspace")
    try:
        df = duckdb_manager.query(f"SELECT * FROM {table_name}")
        out_dir = "exports"
        os.makedirs(out_dir, exist_ok=True)
        out_path = os.path.join(out_dir, f"workspace_{workspace_id}_cleaned.xlsx")
        df.to_excel(out_path, index=False, engine='openpyxl')
        return FileResponse(out_path, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", filename=f"{ws.get('name')}_export.xlsx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

